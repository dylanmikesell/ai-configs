"""Generate references/layout-gallery.pptx: one slide per branded layout.

Each slide uses one of the template's layouts, with the title placeholder set to
"[index] Layout name" and every text placeholder labelled with its idx, so you
can open the file in PowerPoint and see exactly what each branded layout looks
like and which placeholder is which. Picture placeholders are left empty (they
show PowerPoint's image prompt).

Regenerate after the template changes:  python scripts/make_gallery.py
"""
from pathlib import Path

from pptx.enum.shapes import PP_PLACEHOLDER

from pptx_helpers import blank_deck

OUT = Path(__file__).resolve().parent.parent / "references" / "layout-gallery.pptx"

PICTURE_TYPES = {PP_PLACEHOLDER.PICTURE}


def main():
    prs = blank_deck()
    index = 0
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            slide = prs.slides.add_slide(layout)
            for ph in slide.placeholders:
                pf = ph.placeholder_format
                if pf.type in PICTURE_TYPES:
                    continue  # leave picture prompts visible
                if pf.idx == 0:
                    ph.text = f"[{index}] {layout.name}"
                else:
                    ph.text = f"idx {pf.idx}"
            index += 1
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Wrote {index} layout slides -> {OUT}")


if __name__ == "__main__":
    main()
