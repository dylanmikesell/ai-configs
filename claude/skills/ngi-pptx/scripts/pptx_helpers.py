"""Helpers for building NGI-branded presentations from the corporate template.

The whole point of this skill: ALWAYS start from the template file so the custom
slide size (33x19 cm), the master graphics, the Aptos fonts and the NGI colour
palette are *inherited*, never rebuilt from python-pptx defaults. These helpers
make that the path of least resistance.

Usage sketch:

    from pptx_helpers import blank_deck, add_slide, set_title, set_bullets, add_picture

    prs = blank_deck()                                    # template, example slides removed
    cover = add_slide(prs, "Forside #1")
    set_title(cover, "Project kickoff")
    set_text(cover, 1, "NGI · 2026")               # subtitle placeholder (idx 1)

    body = add_slide(prs, "Overskrift og to tekstbokser")  # two-column content
    set_title(body, "Findings")
    set_bullets(body, 14, ["Point A", ("sub-point", 1), "Point B"])
    set_bullets(body, 15, ["Right column"])

    prs.save("out.pptx")

Every layout's placeholder indices are documented in ../references/layouts.md.
"""
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

# assets/ngi-template.pptx sits next to this scripts/ dir, one level up.
TEMPLATE = Path(__file__).resolve().parent.parent / "assets" / "ngi-template.pptx"

# Placeholder types that hold a picture, not text (skip .text on these).
PICTURE_TYPES = {PP_PLACEHOLDER.PICTURE}


def open_template():
    """Open the NGI template with its example slides intact (5 sample slides)."""
    if not TEMPLATE.exists():
        raise FileNotFoundError(f"NGI template not found at {TEMPLATE}")
    return Presentation(str(TEMPLATE))


def blank_deck():
    """Template opened with all example slides removed but ALL branding kept.

    python-pptx has no public 'delete slide' API. We remove each slide id AND
    drop its relationship, so the underlying slide part becomes unreachable and
    is discarded on save. (Removing only the sldId leaves an orphaned part,
    which collides on partname allocation and produces a corrupt file.) Masters,
    layouts, theme and slide size are untouched, so new slides stay branded.
    """
    prs = open_template()
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        prs.part.drop_rel(sldId.rId)
        sldIdLst.remove(sldId)
    return prs


def get_layout(prs, name):
    """Return a slide layout by its EXACT name (see references/layouts.md).

    Names are Norwegian and some have quirks (e.g. ' Tidslinje2' has a leading
    space). Matching by name is more robust than by index.
    """
    available = []
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            available.append(layout.name)
            if layout.name == name:
                return layout
    raise KeyError(
        f"Layout {name!r} not found. Available: "
        + ", ".join(repr(n) for n in available)
    )


def add_slide(prs, layout_name):
    """Add a new slide built from the named branded layout; return the slide."""
    return prs.slides.add_slide(get_layout(prs, layout_name))


def _placeholder(slide, idx):
    try:
        return slide.placeholders[idx]
    except KeyError:
        have = ", ".join(str(p.placeholder_format.idx) for p in slide.placeholders)
        raise KeyError(f"No placeholder idx={idx} on this slide. Present idx: {have}")


def set_title(slide, text):
    """Set the title. In this template the title is a ctrTitle at idx 0."""
    _placeholder(slide, 0).text = text


def set_text(slide, idx, text):
    """Set a single string into the placeholder at ``idx``."""
    _placeholder(slide, idx).text = text


def set_bullets(slide, idx, items):
    """Fill a body placeholder with bullets.

    ``items`` is a list where each entry is either a string (level 0) or a
    ``(text, level)`` tuple. Inherited bullet styling from the layout is kept.
    """
    tf = _placeholder(slide, idx).text_frame
    tf.clear()
    for j, item in enumerate(items):
        text, level = item if isinstance(item, tuple) else (item, 0)
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = text
        p.level = level


def add_picture(slide, idx, image_path):
    """Insert a picture into a picture placeholder (fills + crops to the frame)."""
    return _placeholder(slide, idx).insert_picture(str(image_path))


def describe(prs):
    """Print every layout with its placeholders — handy sanity check in a session."""
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            print(f"# {layout.name!r}")
            for ph in layout.placeholders:
                pf = ph.placeholder_format
                print(f"    idx={pf.idx:<3} type={pf.type} name={ph.name!r}")
